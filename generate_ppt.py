#!/usr/bin/env python3
"""BOSS-2 Presentation Generator — ppt-design system v1"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

C = {
    'primary':      RGBColor(0x5b, 0x8f, 0xd9),
    'primary_deep': RGBColor(0x4a, 0x7b, 0xc0),
    'primary_soft': RGBColor(0xdc, 0xe7, 0xf7),
    'canvas':       RGBColor(0xff, 0xff, 0xff),
    'sky':          RGBColor(0xea, 0xf1, 0xfa),
    'mint':         RGBColor(0xe8, 0xf1, 0xeb),
    'blush':        RGBColor(0xf5, 0xea, 0xe6),
    'sand':         RGBColor(0xf5, 0xef, 0xe6),
    'lavender':     RGBColor(0xef, 0xe9, 0xf5),
    'dark_bg':      RGBColor(0x27, 0x27, 0x29),
    'dark_tile':    RGBColor(0x35, 0x35, 0x3a),
    'dark_tile2':   RGBColor(0x33, 0x38, 0x3a),
    'dark_tile3':   RGBColor(0x35, 0x33, 0x3a),
    'ink':          RGBColor(0x2d, 0x2d, 0x33),
    'ink80':        RGBColor(0x4a, 0x4a, 0x52),
    'ink48':        RGBColor(0x8e, 0x8e, 0x94),
    'white':        RGBColor(0xff, 0xff, 0xff),
    'dark_link':    RGBColor(0x6b, 0xa3, 0xec),
    'dark_body':    RGBColor(0xcc, 0xcc, 0xcc),
    'hairline':     RGBColor(0xe0, 0xe0, 0xe0),
}

FONT = "Noto Sans KR"


def X(px): return Inches(px / SLIDE_W_PX * SLIDE_W_IN)
def Y(px): return Inches(px / SLIDE_H_PX * SLIDE_H_IN)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None):
    s = slide.shapes.add_shape(1, X(x), Y(y), X(w), Y(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, size, bold=False, color=None,
        align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(X(x), Y(y), X(w), Y(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def header(slide, chapter, title, subtitle, page_n):
    txt(slide, chapter, 80, 58, 900, 28, 10, bold=True, color=C['primary'])
    txt(slide, title, 80, 92, 1600, 68, 28, bold=True, color=C['ink'])
    txt(slide, subtitle, 80, 168, 1600, 46, 16, color=C['ink80'])
    txt(slide, f"BOSS-2  ·  {chapter}", 80, 1044, 1400, 22, 9, color=C['ink48'])
    txt(slide, str(page_n), 1820, 1044, 60, 22, 9, color=C['ink48'],
        align=PP_ALIGN.RIGHT)


def chip(slide, x, y, label):
    cw = min(len(label) * 12 + 28, 180)
    rect(slide, x, y, cw, 26, fill=C['white'])
    txt(slide, label, x, y + 2, cw, 22, 10, bold=True, color=C['primary_deep'],
        align=PP_ALIGN.CENTER)
    return 38  # height consumed


def tile(slide, x, y, w, h, bg_c, chip_lbl=None, title=None, body=None):
    rect(slide, x, y, w, h, fill=bg_c)
    cy = y + 16
    if chip_lbl:
        cy += chip(slide, x + 16, cy, chip_lbl)
    if title:
        txt(slide, title, x + 16, cy, w - 32, 44, 15, bold=True, color=C['ink'])
        cy += 46
    if body:
        txt(slide, body, x + 16, cy, w - 32, 80, 12, color=C['ink80'])


def kpi_rail(slide, y, h, items):
    n = len(items)
    gap = 24
    total_gap = gap * (n - 1)
    cell_w = (1760 - total_gap) // n
    for i, (bg_c, val, lbl) in enumerate(items):
        kx = 80 + i * (cell_w + gap)
        rect(slide, kx, y, cell_w, h, fill=bg_c)
        txt(slide, val, kx, y + 14, cell_w, h // 2, 28, bold=True, color=C['ink'],
            align=PP_ALIGN.CENTER)
        txt(slide, lbl, kx, y + h - 36, cell_w, 28, 10, bold=True,
            color=C['primary_deep'], align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════
# S1 — COVER (dark)
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, C['dark_bg'])
txt(s1, "BOSS-2", 160, 310, 1600, 180, 72, bold=True, color=C['white'],
    align=PP_ALIGN.CENTER)
txt(s1, "AI 기반 소상공인 자율 운영 플랫폼", 260, 500, 1400, 60, 24,
    color=C['dark_body'], align=PP_ALIGN.CENTER)
txt(s1, "Planner · Recruitment · Marketing · Sales · Documents",
    360, 578, 1200, 36, 13, color=C['dark_link'], align=PP_ALIGN.CENTER)
# Divider line
rect(s1, 640, 630, 640, 1, fill=RGBColor(0x44, 0x44, 0x48))
txt(s1, "LangGraph StateGraph + DeepAgent 2계층 + Celery Beat",
    560, 644, 800, 28, 10, color=C['ink48'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# S2 — 소상공인의 고충 (light)
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, C['canvas'])
header(s2, "01 · 배경", "소상공인의 운영 고충",
       "채용·마케팅·매출 분석 — 혼자 감당하기엔 너무 많은 업무", 2)

problem_tiles = [
    (C['sky'],   "채용 관리", "34%", "전체 업무의 34%\n채용 관련 행정 처리"),
    (C['mint'],  "마케팅",   "40h", "매달 평균 40시간\nSNS · 광고 운영"),
    (C['blush'], "매출 분석", "52%", "소상공인 52%가\n데이터 미활용"),
]
for i, (c, lbl, stat, sub) in enumerate(problem_tiles):
    tx = 80 + i * 579
    rect(s2, tx, 280, 555, 490, fill=c)
    chip(s2, tx + 16, 296, lbl)
    txt(s2, stat, tx, 360, 555, 110, 56, bold=True, color=C['ink'],
        align=PP_ALIGN.CENTER)
    txt(s2, sub, tx + 20, 486, 515, 64, 13, color=C['ink80'],
        align=PP_ALIGN.CENTER)

txt(s2, "※ 소상공인진흥공단 2024 실태조사 기반 추정치", 80, 802, 900, 26, 9,
    color=C['ink48'])

# ══════════════════════════════════════════════
# S3 — BOSS-2 솔루션 (light)
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3, C['canvas'])
header(s3, "02 · 솔루션", "BOSS-2가 해결합니다",
       "4개 도메인 AI 에이전트가 운영 전반을 자율 처리", 3)

# Left hero stat
rect(s3, 80, 280, 840, 460, fill=C['sky'])
chip(s3, 96, 296, "핵심 가치")
txt(s3, "4개", 80, 358, 840, 120, 62, bold=True, color=C['ink'],
    align=PP_ALIGN.CENTER)
txt(s3, "도메인 에이전트", 80, 480, 840, 46, 22, bold=True, color=C['ink'],
    align=PP_ALIGN.CENTER)
txt(s3, "24/7 자율 운영 · 스케줄 실행 · 장기 기억",
    80, 530, 840, 34, 14, color=C['ink80'], align=PP_ALIGN.CENTER)

# Right column (3 stacked tiles)
right = [
    (C['mint'],     "자율 스케줄링",  "Celery Beat · 토글 1개로 예약 자동 실행"),
    (C['sand'],     "장기 기억",     "Redis 단기 + pgvector 장기 컨텍스트 보존"),
    (C['lavender'], "RAG 지식베이스", "pgvector + FTS 하이브리드 · 4 chunks 기본"),
]
for i, (c, t, b) in enumerate(right):
    tile(s3, 944, 280 + i * 154, 816, 128, c, None, t, b)

# ══════════════════════════════════════════════
# S4 — LangGraph 아키텍처 (light)
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4, C['canvas'])
header(s4, "03 · 아키텍처", "LangGraph 멀티에이전트 구조",
       "Planner가 PlanResult 생성 → 도메인 에이전트 병렬 실행 (Send)", 4)

# Process strip
steps = [
    ("01", "사용자 입력",    "message · history"),
    ("02", "Planner Agent", "계획 수립 · 라우팅"),
    ("03", "Domain Agents", "병렬 실행 Send()"),
    ("04", "Synthesizer",   "결과 통합"),
    ("05", "Profile Saver", "Supabase 저장"),
]
step_w = 316
for i, (num, lbl, cap) in enumerate(steps):
    sx = 80 + i * (step_w + 25)
    is_planner = (i == 1)
    sq_bg = C['primary'] if is_planner else C['primary_soft']
    sq_tx = C['white'] if is_planner else C['primary_deep']
    sqx = sx + (step_w - 48) // 2
    rect(s4, sqx, 302, 48, 48, fill=sq_bg)
    txt(s4, num, sqx, 310, 48, 32, 13, bold=True, color=sq_tx,
        align=PP_ALIGN.CENTER)
    if i < 4:
        rect(s4, sqx + 48, 324, step_w - 48 + 25, 2, fill=C['hairline'])
    txt(s4, lbl, sx, 364, step_w, 30, 13, bold=True, color=C['ink'],
        align=PP_ALIGN.CENTER)
    txt(s4, cap, sx, 396, step_w, 26, 11, color=C['ink80'],
        align=PP_ALIGN.CENTER)

# Detail tiles
detail = [
    (C['sky'],      "BossState",    "account_id · message · plan · domain_results"),
    (C['mint'],     "PlanResult",   "mode · opening · steps · choices · profile_updates"),
    (C['blush'],    "MemorySaver",  "thread_id = account_id 체크포인트"),
    (C['lavender'], "Capability",   "describe() → list[dict]  |  run_as_agent()"),
]
dw = 415
for i, (c, ch, body) in enumerate(detail):
    tile(s4, 80 + i * (dw + 25), 472, dw, 212, c, ch, None, body)

# ══════════════════════════════════════════════
# S5 — 4개 도메인 에이전트 (light)
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5, C['canvas'])
header(s5, "03 · 기술", "4가지 도메인 에이전트",
       "각 에이전트는 capability 시스템으로 동적 확장 가능", 5)

domains = [
    (C['sky'],      "Recruitment",  "4", "공고 작성 · 지원자 평가\n면접 일정 · 합격 통보"),
    (C['mint'],     "Marketing",    "5", "SNS 콘텐츠 · 광고 카피\n마케팅 전략 · 키워드"),
    (C['blush'],    "Sales",        "5", "POS 분석 · 매출 리포트\n비용 추적 · 재고 관리"),
    (C['lavender'], "Documents",    "4", "계약서 · 세무 · 법무\n운영 문서 처리"),
]
dw = 415
for i, (c, ch, num, body) in enumerate(domains):
    dx = 80 + i * (dw + 25)
    rect(s5, dx, 280, dw, 490, fill=c)
    chip(s5, dx + 16, 296, ch)
    txt(s5, num, dx, 360, dw, 110, 56, bold=True, color=C['ink'],
        align=PP_ALIGN.CENTER)
    txt(s5, "capabilities", dx, 474, dw, 28, 13, color=C['primary_deep'],
        align=PP_ALIGN.CENTER)
    txt(s5, body, dx + 20, 514, dw - 40, 80, 12, color=C['ink80'],
        align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# S6 — 자율 스케줄링 (light)
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6, C['canvas'])
header(s6, "04 · 자율화", "Celery Beat 자율 스케줄링",
       "artifact metadata 토글 하나로 자율 실행 — schedule_enabled = true", 6)

# Left hero
rect(s6, 80, 280, 840, 320, fill=C['sky'])
chip(s6, 96, 296, "스케줄 구조")
txt(s6, "toggle 1개", 80, 350, 840, 80, 36, bold=True, color=C['ink'],
    align=PP_ALIGN.CENTER)
txt(s6, "metadata.schedule_enabled = true", 80, 440, 840, 34, 13,
    color=C['primary_deep'], align=PP_ALIGN.CENTER)
txt(s6, "→ Celery Beat가 cron 실행 · KST 기준",
    80, 476, 840, 28, 12, color=C['ink80'], align=PP_ALIGN.CENTER)

# Right process
right6 = [
    (C['mint'],     "알림 시스템", "D-7/3/1/0 사전 알림 → activity_logs"),
    (C['blush'],    "실행 로그",  "kind='log' artifact + logged_from 엣지"),
    (C['lavender'], "실패 감지",  "schedule_status · next_run 추적"),
]
for i, (c, t, b) in enumerate(right6):
    tile(s6, 944, 280 + i * 107, 816, 82, c, None, t, b)

# KPI rail
kpi_rail(s6, 634, 134, [
    (C['sky'],      "60s",  "Celery tick"),
    (C['mint'],     "7일",  "장기 메모리 TTL"),
    (C['blush'],    "20턴", "압축 임계값"),
    (C['lavender'], "4개",  "병렬 실행 도메인"),
])

# ══════════════════════════════════════════════
# S7 — 기술 스택 (light)
# ══════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg(s7, C['canvas'])
header(s7, "05 · 스택", "검증된 기술 스택",
       "프론트엔드부터 AI 파이프라인까지 최신 기술 조합", 7)

stacks = [
    (C['sky'],   "Frontend",  [
        ("Next.js 16",       "App Router · Auth Proxy (proxy.ts)"),
        ("Supabase Realtime","캔버스 즉시 반영 (boss:artifacts-changed)"),
        ("React",            "NodeDetailModal 단일 마운트 전략"),
    ]),
    (C['mint'],  "Backend",   [
        ("FastAPI",          "async 전역 · 27개 라우터"),
        ("LangGraph",        "StateGraph · MemorySaver 체크포인트"),
        ("DeepAgents SDK",   "Planner + 4 Domain 2계층 구조"),
    ]),
    (C['blush'], "Data & AI", [
        ("pgvector + FTS",   "RRF 하이브리드 RAG · 4 chunks"),
        ("Celery Beat",      "KST · enable_utc=False"),
        ("BAAI/bge-m3",      "로컬 1024-dim 임베딩"),
    ]),
]
sw = 555
for i, (c, ch, items) in enumerate(stacks):
    sx = 80 + i * (sw + 24)
    rect(s7, sx, 280, sw, 500, fill=c)
    chip(s7, sx + 16, 296, ch)
    for j, (tech, desc) in enumerate(items):
        iy = 350 + j * 150
        txt(s7, tech, sx + 20, iy, sw - 40, 36, 16, bold=True, color=C['ink'])
        txt(s7, desc, sx + 20, iy + 34, sw - 40, 30, 12, color=C['ink80'])
        if j < 2:
            rect(s7, sx + 20, iy + 76, sw - 40, 1, fill=C['hairline'])

# ══════════════════════════════════════════════
# S8 — 비전 (dark)
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
bg(s8, C['dark_bg'])
txt(s8, "06 · 비전", 80, 58, 900, 28, 10, bold=True, color=C['dark_link'])
txt(s8, "소상공인의 AI 파트너", 80, 92, 1760, 70, 32, bold=True, color=C['white'])
txt(s8, "BOSS-2는 도구가 아닙니다. 함께 일하는 동료입니다.",
    80, 170, 1760, 46, 18, color=C['dark_body'])

vision = [
    (C['dark_tile'],  "자율 운영",
     "사장님이 자리를 비워도\n비즈니스는 계속 돌아갑니다"),
    (C['dark_tile2'], "지속 학습",
     "대화할수록 더 똑똑해지는\n맞춤형 운영 파트너"),
    (C['dark_tile3'], "무한 확장",
     "18개 서브허브 →\n더 많은 도메인으로 성장"),
]
vw = 555
for i, (c, t, b) in enumerate(vision):
    vx = 80 + i * (vw + 24)
    rect(s8, vx, 286, vw, 420, fill=c)
    txt(s8, t, vx + 24, 326, vw - 48, 48, 20, bold=True, color=C['white'])
    txt(s8, b, vx + 24, 386, vw - 48, 80, 14, color=C['dark_body'])

# Divider + closing quote
rect(s8, 560, 744, 800, 1, fill=RGBColor(0x44, 0x44, 0x48))
txt(s8, '"운영은 AI에게, 사장님은 본업에 집중"',
    80, 758, 1760, 54, 18, color=C['dark_link'], align=PP_ALIGN.CENTER)

txt(s8, "BOSS-2  ·  06 · 비전", 80, 1044, 1400, 22, 9, color=C['ink48'])
txt(s8, "8", 1820, 1044, 60, 22, 9, color=C['ink48'], align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
output = "D:/dev/BOSS-2/BOSS2_presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
